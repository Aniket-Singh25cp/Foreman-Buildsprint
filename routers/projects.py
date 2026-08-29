import os
import json
import uuid
import tempfile
from typing import Dict
import httpx
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from models import ProjectCreate, ProjectResponse, ProjectData, Task, TaskUpdate

router = APIRouter(prefix="/projects", tags=["projects"])

# In-memory database storing project_id -> ProjectData
PROJECTS_DB: Dict[str, ProjectData] = {}

GATEWAY_URL = "https://latentstack.dev/v1/chat/completions"
MODEL_NAME = "gemini-3.6-flash"

SYSTEM_PROMPT = """You are an expert AI software architect and project manager.

Your job is twofold:
1. Validate whether the user's input is a coherent, genuine project/product brief (even a short one), versus gibberish, keyboard mashing, or a single unrelated word (e.g., "app", "test", "asdfghjk").
2. If it is NOT a valid brief, classify it as invalid and provide a clear, helpful reason explaining why and how to write a valid brief.
3. If it IS a valid brief, decompose it into 3 to 6 structured engineering tasks that map to real product development phases (e.g., Requirements/Scope, UI/UX Design, Backend/API, Frontend/Implementation, Testing & QA, Deployment - adapting naturally to what the brief actually needs).

You MUST return ONLY valid raw JSON with no Markdown formatting, code blocks, or preamble.
The JSON must strictly conform to ONE of these two schema shapes:

If the input is INVALID:
{
  "valid": false,
  "reason": "This doesn't look like a project brief. Try describing what you want to build, like 'a mobile app for tracking gym workouts'."
}

If the input IS VALID:
{
  "valid": true,
  "tasks": [
    {
      "id": 1,
      "title": "Requirements & Scope Definition",
      "description": "Detailed description of engineering work",
      "depends_on": [],
      "status": "pending"
    }
  ]
}

Rules for tasks:
- Tasks should map to recognizable real engineering phases and read like actionable engineering work.
- id: integer starting from 1 sequentially.
- title: concise, phase-aligned engineering task title.
- description: clear, actionable explanation of what needs to be built.
- depends_on: list of task integer IDs that must precede this task (e.g. Frontend depends on UI/UX and Backend).
- status: string, always "pending".
"""

@router.post("", response_model=ProjectResponse)
def create_project(project: ProjectCreate):
    trimmed_brief = project.brief.strip()
    if len(trimmed_brief) < 15:
        raise HTTPException(
            status_code=400,
            detail="Brief is too short - please describe your project in more detail"
        )

    api_key = os.getenv("LATENTSTACK_API_KEY")
    if not api_key:
        raise HTTPException(
            status_code=502,
            detail="LATENTSTACK_API_KEY environment variable is not set."
        )

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }

    payload = {
        "model": MODEL_NAME,
        "messages": [
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": f"Project Brief: {project.brief}"}
        ],
        "temperature": 0.2
    }

    try:
        response = httpx.post(GATEWAY_URL, headers=headers, json=payload, timeout=30.0)
    except httpx.RequestError as exc:
        raise HTTPException(
            status_code=502,
            detail=f"Failed to communicate with LLM gateway: {str(exc)}"
        )

    if response.status_code != 200:
        raise HTTPException(
            status_code=502,
            detail=f"LLM gateway returned error status {response.status_code}: {response.text}"
        )

    try:
        data = response.json()
        content = data["choices"][0]["message"]["content"].strip()
        
        # Clean potential markdown code fences if model outputted ```json ... ```
        if content.startswith("```"):
            lines = content.splitlines()
            if lines[0].startswith("```"):
                lines = lines[1:]
            if lines and lines[-1].startswith("```"):
                lines = lines[:-1]
            content = "\n".join(lines).strip()

        parsed_json = json.loads(content)

        if not isinstance(parsed_json, dict):
            raise ValueError("Expected JSON response to be an object.")

        is_valid = parsed_json.get("valid", True)
        if not is_valid:
            reason = parsed_json.get(
                "reason",
                "This doesn't look like a project brief. Try describing what you want to build in more detail."
            )
            raise HTTPException(status_code=400, detail=reason)

        raw_tasks = parsed_json.get("tasks")
        if not isinstance(raw_tasks, list):
            raise ValueError("Expected 'tasks' field in JSON response to be a list.")

        tasks = [Task(**task_data) for task_data in raw_tasks]
    except HTTPException:
        raise
    except (json.JSONDecodeError, KeyError, ValueError, TypeError) as exc:
        raise HTTPException(
            status_code=502,
            detail=f"Failed to parse model JSON response: {str(exc)}. Raw response: {response.text}"
        )

    project_id = str(uuid.uuid4())
    project_data = ProjectData(project_id=project_id, brief=project.brief, tasks=tasks)
    PROJECTS_DB[project_id] = project_data

    return ProjectResponse(project_id=project_id, tasks=tasks)


@router.get("/{project_id}/summary-deck")
def generate_summary_deck(project_id: str):
    if project_id not in PROJECTS_DB:
        raise HTTPException(status_code=404, detail="Project not found")

    project_data = PROJECTS_DB[project_id]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    INK_BLUE = RGBColor(12, 28, 56)
    AMBER_ACCENT = RGBColor(217, 119, 6)
    EMERALD_ACCENT = RGBColor(16, 185, 129)
    TEXT_DARK = RGBColor(30, 41, 59)
    TEXT_MUTED = RGBColor(100, 116, 139)

    blank_layout = prs.slide_layouts[6]

    # --- SLIDE 1: Title Slide ---
    slide1 = prs.slides.add_slide(blank_layout)

    # Title box
    tb1_title = slide1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(1.0))
    tf1_title = tb1_title.text_frame
    tf1_title.word_wrap = True
    p1 = tf1_title.paragraphs[0]
    p1.text = "PROJECT SUMMARY DECK"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = INK_BLUE

    # Brief box
    tb1_brief = slide1.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(4.0))
    tf1_brief = tb1_brief.text_frame
    tf1_brief.word_wrap = True

    p1_brief_hdr = tf1_brief.paragraphs[0]
    p1_brief_hdr.text = "PROJECT BRIEF:"
    p1_brief_hdr.font.size = Pt(18)
    p1_brief_hdr.font.bold = True
    p1_brief_hdr.font.color.rgb = AMBER_ACCENT

    p1_brief_txt = tf1_brief.add_paragraph()
    p1_brief_txt.text = project_data.brief
    p1_brief_txt.font.size = Pt(20)
    p1_brief_txt.font.color.rgb = TEXT_DARK
    p1_brief_txt.space_before = Pt(12)

    # --- SLIDE 2: Task Status Breakdown ---
    slide2 = prs.slides.add_slide(blank_layout)

    tb2_title = slide2.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(0.8))
    tf2_title = tb2_title.text_frame
    p2 = tf2_title.paragraphs[0]
    p2.text = "TASK BREAKDOWN"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = INK_BLUE

    pending_tasks = [t for t in project_data.tasks if t.status == "pending"]
    done_tasks = [t for t in project_data.tasks if t.status == "done"]

    # Pending Column (Left)
    tb2_pending = slide2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.4), Inches(5.0))
    tf2_pending = tb2_pending.text_frame
    tf2_pending.word_wrap = True

    p2_p_hdr = tf2_pending.paragraphs[0]
    p2_p_hdr.text = f"PENDING ({len(pending_tasks)})"
    p2_p_hdr.font.size = Pt(18)
    p2_p_hdr.font.bold = True
    p2_p_hdr.font.color.rgb = AMBER_ACCENT

    if pending_tasks:
        for t in pending_tasks:
            pt = tf2_pending.add_paragraph()
            pt.text = f"• #{t.id}: {t.title}"
            pt.font.size = Pt(15)
            pt.font.bold = True
            pt.font.color.rgb = TEXT_DARK
            pt.space_before = Pt(8)

            if t.description:
                pd = tf2_pending.add_paragraph()
                pd.text = f"  {t.description}"
                pd.font.size = Pt(12)
                pd.font.color.rgb = TEXT_MUTED
    else:
        pt = tf2_pending.add_paragraph()
        pt.text = "None"
        pt.font.size = Pt(14)
        pt.font.color.rgb = TEXT_MUTED
        pt.space_before = Pt(8)

    # Done Column (Right)
    tb2_done = slide2.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(5.4), Inches(5.0))
    tf2_done = tb2_done.text_frame
    tf2_done.word_wrap = True

    p2_d_hdr = tf2_done.paragraphs[0]
    p2_d_hdr.text = f"COMPLETED ({len(done_tasks)})"
    p2_d_hdr.font.size = Pt(18)
    p2_d_hdr.font.bold = True
    p2_d_hdr.font.color.rgb = EMERALD_ACCENT

    if done_tasks:
        for t in done_tasks:
            dt = tf2_done.add_paragraph()
            dt.text = f"• #{t.id}: {t.title}"
            dt.font.size = Pt(15)
            dt.font.bold = True
            dt.font.color.rgb = TEXT_DARK
            dt.space_before = Pt(8)

            if t.description:
                dd = tf2_done.add_paragraph()
                dd.text = f"  {t.description}"
                dd.font.size = Pt(12)
                dd.font.color.rgb = TEXT_MUTED
    else:
        dt = tf2_done.add_paragraph()
        dt.text = "None"
        dt.font.size = Pt(14)
        dt.font.color.rgb = TEXT_MUTED
        dt.space_before = Pt(8)

    # --- SLIDE 3: Stats Summary ---
    slide3 = prs.slides.add_slide(blank_layout)

    tb3_title = slide3.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(0.8))
    tf3_title = tb3_title.text_frame
    p3 = tf3_title.paragraphs[0]
    p3.text = "PROJECT METRICS"
    p3.font.size = Pt(32)
    p3.font.bold = True
    p3.font.color.rgb = INK_BLUE

    total_count = len(project_data.tasks)
    completed_count = len(done_tasks)
    completion_rate = (completed_count / total_count * 100) if total_count > 0 else 0.0

    # Stat Card 1: Total Tasks
    tb3_s1 = slide3.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(3.5), Inches(3.0))
    tf3_s1 = tb3_s1.text_frame
    p3_s1_v = tf3_s1.paragraphs[0]
    p3_s1_v.text = str(total_count)
    p3_s1_v.font.size = Pt(64)
    p3_s1_v.font.bold = True
    p3_s1_v.font.color.rgb = INK_BLUE

    p3_s1_l = tf3_s1.add_paragraph()
    p3_s1_l.text = "TOTAL TASKS"
    p3_s1_l.font.size = Pt(16)
    p3_s1_l.font.bold = True
    p3_s1_l.font.color.rgb = TEXT_MUTED

    # Stat Card 2: Completed Tasks
    tb3_s2 = slide3.shapes.add_textbox(Inches(4.9), Inches(2.5), Inches(3.5), Inches(3.0))
    tf3_s2 = tb3_s2.text_frame
    p3_s2_v = tf3_s2.paragraphs[0]
    p3_s2_v.text = str(completed_count)
    p3_s2_v.font.size = Pt(64)
    p3_s2_v.font.bold = True
    p3_s2_v.font.color.rgb = EMERALD_ACCENT

    p3_s2_l = tf3_s2.add_paragraph()
    p3_s2_l.text = "COMPLETED"
    p3_s2_l.font.size = Pt(16)
    p3_s2_l.font.bold = True
    p3_s2_l.font.color.rgb = TEXT_MUTED

    # Stat Card 3: Completion Rate
    tb3_s3 = slide3.shapes.add_textbox(Inches(8.8), Inches(2.5), Inches(3.5), Inches(3.0))
    tf3_s3 = tb3_s3.text_frame
    p3_s3_v = tf3_s3.paragraphs[0]
    p3_s3_v.text = f"{completion_rate:.0f}%"
    p3_s3_v.font.size = Pt(64)
    p3_s3_v.font.bold = True
    p3_s3_v.font.color.rgb = AMBER_ACCENT

    p3_s3_l = tf3_s3.add_paragraph()
    p3_s3_l.text = "COMPLETION RATE"
    p3_s3_l.font.size = Pt(16)
    p3_s3_l.font.bold = True
    p3_s3_l.font.color.rgb = TEXT_MUTED

    temp_dir = tempfile.gettempdir()
    file_path = os.path.join(temp_dir, f"summary_deck_{project_id}.pptx")
    prs.save(file_path)

    media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    filename = f"project_{project_id}_summary.pptx"

    return FileResponse(
        path=file_path,
        media_type=media_type,
        filename=filename
    )


@router.patch("/{project_id}/tasks/{task_id}", response_model=Task)
def update_task_status(project_id: str, task_id: int, task_update: TaskUpdate):
    if project_id not in PROJECTS_DB:
        raise HTTPException(status_code=404, detail="Project not found")

    project = PROJECTS_DB[project_id]
    target_task = None
    for task in project.tasks:
        if task.id == task_id:
            target_task = task
            break

    if not target_task:
        raise HTTPException(status_code=404, detail="Task not found")

    target_task.status = task_update.status
    return target_task


